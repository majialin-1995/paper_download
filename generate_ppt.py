#!/usr/bin/env python3
"""
根据 summaries 目录生成演示文稿（占位符替换版）
"""

from __future__ import annotations

import argparse
import copy
import json
import re
import shutil
from pathlib import Path
from typing import Any, Dict, List

try:
    from pptx import Presentation
except Exception:
    Presentation = None


# ───────────────────────── 参数解析 ───────────────────────── #

def parse_args(argv: List[str] | None = None) -> argparse.Namespace:
    p = argparse.ArgumentParser("根据 summaries 生成 PPT（占位符版）")
    p.add_argument("summaries", type=Path, help="包含 JSON 摘要的目录")
    p.add_argument(
        "--template", type=Path, default=Path("template.pptx"),
        help="PPT 模板文件路径"
    )
    p.add_argument(
        "--out", type=Path, default=Path("slides.pptx"),
        help="输出 PPT 文件名"
    )
    p.add_argument(
        "--refs", type=Path, default=None,
        help="参考文献文本（每行一条）。默认：summaries 上级目录的 references_ieee.txt"
    )
    p.add_argument(
        "--print-info", action="store_true",
        help="仅打印摘要信息，不生成 PPT"
    )

    args = p.parse_args(argv)

    if args.refs is None:
        args.refs = args.summaries.parent / "references_ieee.txt"

    return args


# ─────────────────────── 工具函数 ─────────────────────── #

def normalize_text(s: str) -> str:
    return re.sub(r"[^\w\s]", "", s).lower()


def find_reference(title: str, refs: List[str]) -> str:
    """在 refs 里用标题做简单模糊匹配，返回对应那一行参考文献"""
    norm = normalize_text(title)
    for line in refs:
        if norm in normalize_text(line):
            return line
    return ""


def duplicate_slide(prs, slide):
    """复制一页 slide（包括所有 shapes）"""
    new_slide = prs.slides.add_slide(slide.slide_layout)
    for shp in slide.shapes:
        new_el = copy.deepcopy(shp.element)
        new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')
    return new_slide


# ①②③… 生成器
CIRCLES = ["①","②","③","④","⑤","⑥","⑦","⑧","⑨","⑩",
           "⑪","⑫","⑬","⑭","⑮","⑯","⑰","⑱","⑲","⑳"]

def indexed_text(lst: List[Any] | None) -> str:
    if not lst:
        return ""
    out = []
    for i, item in enumerate(lst):
        mark = CIRCLES[i] if i < len(CIRCLES) else f"({i+1})"
        out.append(f"{mark} {item}")
    return "\n".join(out)


def plain_text(lst: List[Any] | None) -> str:
    if not lst:
        return ""
    return "\n".join(str(x) for x in lst)


def copy_pdf_for_json(
    idx: int,
    jf: Path,
    data: Dict[str, Any],
    out_dir: Path,
) -> None:
    """
    从 data['pdf_path'] 拿到原始 PDF 路径，复制到 out_dir 下，
    文件名加上编号前缀，例如 001_xxx.pdf。
    """
    pdf_path = data.get("pdf_path")
    if not pdf_path:
        print(f"⚠ JSON {jf.name} 未包含 pdf_path 字段，跳过复制")
        return

    src = Path(pdf_path)
    # 相对路径相对于 summaries 目录
    if not src.is_absolute():
        src = jf.parent / src

    if not src.is_file():
        print(f"⚠ 未找到对应 PDF（按 pdf_path）：{pdf_path}")
        return

    out_dir.mkdir(parents=True, exist_ok=True)
    dst_name = f"[{idx:02d}]_{src.name}"
    dst = out_dir / dst_name

    try:
        shutil.copy2(src, dst)
        # 可选：打印一下
        print(f"✔ 复制 PDF: {src} -> {dst}")
    except Exception as e:
        print(f"⚠ 复制 PDF 失败（{src} -> {dst}）：{e}")


# ─────────────────────── 占位符填充 ─────────────────────── #

def fill_placeholders(
    slide,
    data: Dict[str, Any],
    title: str,
    reference: str,
    page: int,
    total_pages: int,
    section: str,
    idx: int
) -> None:

    mapping: Dict[str, str] = {
        "{{No.}}":        str(idx),
        "{{title}}":      title,
        "{{reference}}":  reference or "",
        "{{Pages}}":      f"{page} / {total_pages}",
        "{{totalpages}}": str(total_pages),
    }

    if section == "intro":
        mapping["{{phenomenon}}"] = plain_text(data.get("phenomenon") or [])
        mapping["{{problems}}"]   = indexed_text(data.get("problem") or [])
        mapping["{{methods}}"]    = indexed_text(data.get("mechanism") or [])

    elif section == "conclusion":
        mapping["{{results}}"]    = indexed_text(data.get("result") or [])

    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                t = run.text
                for ph, val in mapping.items():
                    if ph in t:
                        t = t.replace(ph, val)
                run.text = t


# ─────────────────────── Overlist 导出 ─────────────────────── #

def export_overlist(used_refs: List[str], refs_file: Path) -> None:
    """
    根据本次实际用到的参考文献，生成 references_overlist.txt：
      [1] ...
      [2] ...
    """
    if not used_refs:
        print("ℹ 无任何引用匹配，跳过 overlist")
        return

    # 去重保持顺序
    seen = set()
    ordered: List[str] = []
    for r in used_refs:
        if not r:
            continue
        if r in seen:
            continue
        seen.add(r)
        ordered.append(r)

    if not ordered:
        print("ℹ 用到的引用全是空字符串，跳过 overlist")
        return

    out_txt = refs_file.parent / "references_overlist.txt"
    with out_txt.open("w", encoding="utf-8") as f:
        for i, ref in enumerate(ordered, 1):
            f.write(f"[{i}] {ref}\n")

    print(f"✔ 已生成 references_overlist.txt: {out_txt}")


# ───────────────────────── 主逻辑 ───────────────────────── #

def main(argv: List[str] | None = None) -> None:
    args = parse_args(argv)

    json_files = sorted(args.summaries.glob("*.json"))
    if not json_files:
        raise SystemExit(f"❌ 未在 {args.summaries} 找到 .json 文件")

    all_refs = (
        args.refs.read_text(encoding="utf-8").splitlines()
        if args.refs.is_file()
        else []
    )

    used_refs: List[str] = []

    # 所有 PDF 最终集中复制到这里
    pdf_out_dir = args.summaries.parent / "pdf_with_idx"

    # 只打印摘要信息
    if args.print_info:
        for i, jf in enumerate(json_files, 1):
            data = json.loads(jf.read_text("utf-8"))
            title = jf.stem.split("_", 1)[1] if "_" in jf.stem else jf.stem
            ref = find_reference(title, all_refs)
            if ref:
                used_refs.append(ref)

            # 复制 PDF（带编号）
            copy_pdf_for_json(i, jf, data, pdf_out_dir)

            print(f"{i}. {title}")
            print(ref)
            print(json.dumps(data, ensure_ascii=False, indent=2), "\n")

        export_overlist(used_refs, args.refs)
        return

    if Presentation is None:
        raise SystemExit("请安装 python-pptx： pip install python-pptx")

    prs = Presentation(str(args.template))

    base_title = base_intro = base_conclusion = None

    for s in prs.slides:
        txt = "\n".join(sh.text for sh in s.shapes if sh.has_text_frame)
        if "{{title}}" in txt and base_title is None:
            base_title = s
        if "{{problems}}" in txt and "{{methods}}" in txt:
            base_intro = s
        if "{{results}}" in txt:
            base_conclusion = s

    if not all([base_title, base_intro, base_conclusion]):
        raise SystemExit("❌ 模板缺失 title/intro/conclusion 页")

    template_slide_count = len(prs.slides)
    extra = 3
    total_pages = template_slide_count + extra * len(json_files)

    for idx, jf in enumerate(json_files, 1):
        data = json.loads(jf.read_text("utf-8"))
        title = jf.stem.split("_", 1)[1] if "_" in jf.stem else jf.stem
        ref = find_reference(title, all_refs)
        if ref:
            used_refs.append(ref)

        # 复制 PDF（带编号）
        copy_pdf_for_json(idx, jf, data, pdf_out_dir)

        page1 = template_slide_count + extra * (idx - 1) + 1
        page2 = page1 + 1
        page3 = page1 + 2

        s1 = duplicate_slide(prs, base_title)
        fill_placeholders(s1, data, title, ref, page1, total_pages, "title", idx)

        s2 = duplicate_slide(prs, base_intro)
        fill_placeholders(s2, data, title, ref, page2, total_pages, "intro", idx)

        s3 = duplicate_slide(prs, base_conclusion)
        fill_placeholders(s3, data, title, ref, page3, total_pages, "conclusion", idx)

    prs.save(args.out)
    print(f"✔ 已生成 PPT: {args.out}")

    # 只生成 overlist，不再处理 PDF
    export_overlist(used_refs, args.refs)


if __name__ == "__main__":
    main()
